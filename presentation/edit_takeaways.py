"""Rewrite slide 9 (TAKEAWAYS) as a simple list."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

CRIMSON = RGBColor(0x98, 0x1E, 0x32)
DEEP    = RGBColor(0x2B, 0x2D, 0x42)
CREAM   = RGBColor(0xF8, 0xF4, 0xF0)
TEAL    = RGBColor(0x0E, 0x7C, 0x86)
GOLD    = RGBColor(0xC5, 0x8E, 0x2A)
MUTED   = RGBColor(0x6B, 0x6B, 0x6B)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

PATH = "/Users/andrewedson/Desktop/CPTS 440/Virtual Counsler/Virtual_Counselor_CPTS440.pptx"
W, H = Inches(13.3), Inches(7.5)

prs = Presentation(PATH)
slide = prs.slides[8]  # slide 9

# nuke existing shapes
spTree = slide.shapes._spTree
for shp in list(slide.shapes):
    spTree.remove(shp._element)

bg = slide.background
bg.fill.solid()
bg.fill.fore_color.rgb = CREAM

def add_rect(x, y, w, h, rgb, line_rgb=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = rgb
    if line_rgb is None: s.line.fill.background()
    else: s.line.color.rgb = line_rgb
    s.shadow.inherit = False
    return s

def add_text(x, y, w, h, text, *, font="Calibri", size=12, color=DEEP,
             bold=False, italic=False, align=PP_ALIGN.LEFT,
             valign=MSO_ANCHOR.TOP, char_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = valign
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        if char_spacing is not None:
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(int(char_spacing * 100)))
    return tb

# crimson side band
add_rect(Inches(0), Inches(0), Inches(0.25), H, CRIMSON)

# title
add_text(Inches(0.7), Inches(0.4), Inches(12), Inches(0.35),
         "TAKEAWAYS", font="Consolas", size=12, color=CRIMSON,
         bold=True, char_spacing=4)
add_text(Inches(0.7), Inches(0.75), Inches(12), Inches(0.85),
         "what we'd tell the next team", font="Georgia", size=32, color=DEEP, bold=True)

# list of takeaways
items = [
    ("use deterministic search when possible",
     "vector search is for fuzzy questions. for known things — a course code, a degree name — just look it up."),
    ("two stages beat one big model",
     "cheap cosine narrows 3,500 chunks down to 30. the reranker only spends compute where it matters."),
    ("let the LLM write, not decide",
     "the rule engine decides if you can take the class. claude writes the explanation. don't swap the two."),
    ("model your data before tuning prompts",
     "DNF for prereqs (AND-of-ORs) was the single biggest accuracy win. no prompt change came close."),
    ("partials still teach you something",
     "75% of our \"useful\" answers come from cases where the model named most of the right courses, not all. that's still a working advisor."),
]

# bullet list — clean, left-aligned, generous spacing
y = Inches(2.0)
for i, (head, body) in enumerate(items, 1):
    # number + heading line
    add_text(Inches(0.9), y, Inches(0.55), Inches(0.45),
             f"{i:02d}", font="Georgia", size=22, color=CRIMSON, bold=True)
    add_text(Inches(1.5), y, Inches(11.2), Inches(0.45),
             head, font="Georgia", size=20, color=DEEP, bold=True)
    # body line
    add_text(Inches(1.5), y + Inches(0.5), Inches(11.2), Inches(0.45),
             body, font="Calibri", size=14, color=MUTED, italic=True)
    y += Inches(0.94)

# footer
add_rect(Inches(0), H - Inches(0.35), W, Inches(0.35), DEEP)
add_text(Inches(0.5), H - Inches(0.35), Inches(8), Inches(0.35),
         "Virtual Counselor  ·  CPTS 440  ·  Group 6",
         font="Calibri", size=10, color=WHITE, valign=MSO_ANCHOR.MIDDLE)
add_text(W - Inches(0.7), H - Inches(0.35), Inches(0.4), Inches(0.35),
         "9", font="Calibri", size=10, color=WHITE,
         valign=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

prs.save(PATH)
print("OK · slide 9 simplified to list format")
