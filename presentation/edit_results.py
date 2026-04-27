"""
Edit slide 8 (RESULTS) of the user's existing pptx in place.
Preserves all of their other edits — only touches slide 8.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from copy import deepcopy
from lxml import etree

# colors
CRIMSON = RGBColor(0x98, 0x1E, 0x32)
DEEP    = RGBColor(0x2B, 0x2D, 0x42)
CREAM   = RGBColor(0xF8, 0xF4, 0xF0)
SAND    = RGBColor(0xED, 0xE6, 0xDD)
TEAL    = RGBColor(0x0E, 0x7C, 0x86)
GOLD    = RGBColor(0xC5, 0x8E, 0x2A)
MUTED   = RGBColor(0x6B, 0x6B, 0x6B)
LINE    = RGBColor(0xC9, 0xBF, 0xB3)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

W, H = Inches(13.3), Inches(7.5)

PATH = "/Users/andrewedson/Desktop/CPTS 440/Virtual Counsler/Virtual_Counselor_CPTS440.pptx"

prs = Presentation(PATH)
slide = prs.slides[7]  # slide 8

# ---- nuke existing shapes on slide 8 ----
spTree = slide.shapes._spTree
for shp in list(slide.shapes):
    spTree.remove(shp._element)

# ---- background fill: cream ----
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = CREAM

# ---- helpers ----
def add_rect(x, y, w, h, rgb, line_rgb=None, line_w=0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = rgb
    if line_rgb is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_rgb
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s

def add_text(x, y, w, h, text, *, font="Calibri", size=12, color=DEEP,
             bold=False, italic=False, align=PP_ALIGN.LEFT,
             valign=MSO_ANCHOR.TOP, char_spacing=None, margin=0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(margin)
    tf.word_wrap = True
    tf.vertical_anchor = valign
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        if char_spacing is not None:
            # spc attribute on rPr (in hundredths of a point)
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(int(char_spacing * 100)))
    return tb

# ---- crimson side band ----
add_rect(Inches(0), Inches(0), Inches(0.25), H, CRIMSON)

# ---- title ----
add_text(Inches(0.7), Inches(0.4), Inches(12), Inches(0.35),
         "RESULTS", font="Consolas", size=12, color=CRIMSON,
         bold=True, char_spacing=4)
add_text(Inches(0.7), Inches(0.75), Inches(12), Inches(0.85),
         "120 cases auto-graded  ·  91 / 120 useful  ·  29 outright fails",
         font="Georgia", size=28, color=DEEP, bold=True)

# ---- top stat strip: 4 boxes ----
def stat_card(x, big, label, color):
    add_rect(x, Inches(1.85), Inches(2.95), Inches(1.45), WHITE, line_rgb=LINE, line_w=0.75)
    add_rect(x, Inches(1.85), Inches(0.08), Inches(1.45), color)
    add_text(x + Inches(0.2), Inches(1.92), Inches(2.7), Inches(0.85),
             big, font="Georgia", size=44, color=DEEP, bold=True)
    add_text(x + Inches(0.2), Inches(2.82), Inches(2.7), Inches(0.45),
             label, font="Calibri", size=11, color=MUTED)

def big_stat(x, big, label, color):
    add_rect(x, Inches(1.85), Inches(2.95), Inches(1.45), WHITE, line_rgb=LINE, line_w=0.75)
    add_rect(x, Inches(1.85), Inches(0.08), Inches(1.45), color)
    add_text(x + Inches(0.2), Inches(1.92), Inches(2.7), Inches(0.85),
             big, font="Georgia", size=34, color=DEEP, bold=True)
    add_text(x + Inches(0.2), Inches(2.82), Inches(2.7), Inches(0.45),
             label, font="Calibri", size=11, color=MUTED)

big_stat(Inches(0.7),    "91 / 120",  "useful answers (pass + partial)",  TEAL)
stat_card(Inches(3.85),  "61",        "pass · exact match",                TEAL)
stat_card(Inches(7.0),   "30",        "partial · mostly right",            GOLD)
stat_card(Inches(10.15), "29",        "fail · wrong answer",               CRIMSON)

# ---- per-category stacked bar chart ----
add_text(Inches(0.7), Inches(3.55), Inches(8), Inches(0.35),
         "by category  ·  pass / partial / fail", font="Consolas",
         size=11, color=DEEP, bold=True, char_spacing=3)

chart_data = CategoryChartData()
chart_data.categories = [
    "prereq valid.", "chain discovery", "credits",
    "schedule", "ucore", "degree progress"
]
# categories order: prereq valid · chain · credits · schedule · ucore · degree
chart_data.add_series("pass",    (20,  4, 11, 13, 7, 6))
chart_data.add_series("partial", ( 1, 18,  2,  1, 3, 5))
chart_data.add_series("fail",    ( 9,  3,  7,  1, 0, 9))

cx, cy, cw, ch = Inches(0.7), Inches(3.95), Inches(7.6), Inches(2.95)
chart_shape = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_STACKED, cx, cy, cw, ch, chart_data
)
chart = chart_shape.chart
chart.has_title = False
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.TOP
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(10)
chart.legend.font.name = "Calibri"

# series colors
series_colors = [TEAL, GOLD, CRIMSON]
for ser, c in zip(chart.series, series_colors):
    fill = ser.format.fill
    fill.solid()
    fill.fore_color.rgb = c
    ser.format.line.fill.background()
    # show value labels on each segment
    dl = ser.data_labels
    dl.show_value = True
    dl.font.size = Pt(9)
    dl.font.name = "Calibri"
    dl.font.color.rgb = WHITE
    dl.font.bold = True

# axis styling
cat_ax = chart.category_axis
cat_ax.tick_labels.font.size = Pt(10)
cat_ax.tick_labels.font.name = "Calibri"
cat_ax.tick_labels.font.color.rgb = MUTED

val_ax = chart.value_axis
val_ax.tick_labels.font.size = Pt(9)
val_ax.tick_labels.font.color.rgb = MUTED
val_ax.maximum_scale = 30
val_ax.minimum_scale = 0

# ---- right card: best/worst categories ----
card_x, card_y, card_w, card_h = Inches(8.5), Inches(3.55), Inches(4.2), Inches(3.35)
add_rect(card_x, card_y, card_w, card_h, DEEP)
add_text(card_x + Inches(0.2), card_y + Inches(0.15), card_w - Inches(0.4), Inches(0.35),
         "where we win and lose", font="Consolas", size=11,
         color=RGBColor(0x9D, 0xD8, 0xDC), bold=True, char_spacing=3)

# best
add_text(card_x + Inches(0.2), card_y + Inches(0.6), card_w - Inches(0.4), Inches(0.32),
         "BEST  ·  ucore planning", font="Consolas", size=10,
         color=RGBColor(0x9D, 0xD8, 0xDC), bold=True, char_spacing=3)
add_text(card_x + Inches(0.2), card_y + Inches(0.92), card_w - Inches(0.4), Inches(0.4),
         "10 / 10 useful  ·  zero outright fails",
         font="Calibri", size=12, color=WHITE)

# strong
add_text(card_x + Inches(0.2), card_y + Inches(1.4), card_w - Inches(0.4), Inches(0.32),
         "STRONG  ·  schedule feasibility", font="Consolas", size=10,
         color=RGBColor(0x9D, 0xD8, 0xDC), bold=True, char_spacing=3)
add_text(card_x + Inches(0.2), card_y + Inches(1.72), card_w - Inches(0.4), Inches(0.4),
         "13 / 15 pass  ·  DNF prereq engine pays off",
         font="Calibri", size=12, color=WHITE)

# weakest
add_text(card_x + Inches(0.2), card_y + Inches(2.2), card_w - Inches(0.4), Inches(0.32),
         "WEAKEST  ·  degree progress", font="Consolas", size=10,
         color=RGBColor(0xFF, 0xE9, 0xA8), bold=True, char_spacing=3)
add_text(card_x + Inches(0.2), card_y + Inches(2.52), card_w - Inches(0.4), Inches(0.7),
         "6 pass / 5 partial / 9 fail  ·  many degrees not in catalog text in clean form",
         font="Calibri", size=12, color=WHITE)

# ---- footer ----
add_rect(Inches(0), H - Inches(0.35), W, Inches(0.35), DEEP)
add_text(Inches(0.5), H - Inches(0.35), Inches(8), Inches(0.35),
         "Virtual Counselor  ·  CPTS 440  ·  Group 6",
         font="Calibri", size=10, color=WHITE, valign=MSO_ANCHOR.MIDDLE)
add_text(W - Inches(0.7), H - Inches(0.35), Inches(0.4), Inches(0.35),
         "8", font="Calibri", size=10, color=WHITE,
         valign=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

prs.save(PATH)
print("OK · slide 8 rebuilt with test-case results")
